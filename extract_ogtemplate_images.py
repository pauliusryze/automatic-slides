from pptx import Presentation

prs = Presentation('OGtemplate.pptx')
slide = prs.slides[0]
print(f'Slide width: {prs.slide_width.inches:.2f} in, height: {prs.slide_height.inches:.2f} in')
print('--- Image Slots (Pictures Only) ---')
for i, shape in enumerate(slide.shapes):
    if shape.shape_type == 13:  # PICTURE
        print(f'[{i}] Name: {getattr(shape, "name", "No name")}')
        print(f'    Position: left={shape.left.inches:.2f}, top={shape.top.inches:.2f}')
        print(f'    Size: width={shape.width.inches:.2f}, height={shape.height.inches:.2f}') 