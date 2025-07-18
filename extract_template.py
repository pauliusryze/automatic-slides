#!/usr/bin/env python3
"""
Extract exact positioning and sizing from template PPTX file.
"""

from pptx import Presentation
from pptx.util import Inches
import os

def extract_template_measurements(template_path):
    """Extract exact positioning and sizing from template."""
    
    if not os.path.exists(template_path):
        print(f"Template file not found: {template_path}")
        return None
    
    prs = Presentation(template_path)
    
    if len(prs.slides) == 0:
        print("No slides found in template")
        return None
    
    slide = prs.slides[0]
    measurements = {
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
        'shapes': []
    }
    
    for i, shape in enumerate(slide.shapes):
        shape_info = {
            'index': i,
            'name': shape.name if hasattr(shape, 'name') else f'Shape_{i}',
            'type': type(shape).__name__,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height
        }
        
        # Extract text if it's a text box
        if hasattr(shape, 'text_frame'):
            shape_info['text'] = shape.text_frame.text
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                shape_info['font_size'] = para.font.size
                shape_info['font_bold'] = para.font.bold
                shape_info['font_color'] = para.font.color.rgb if para.font.color.rgb else None
        
        measurements['shapes'].append(shape_info)
        print(f"Shape {i}: {shape_info}")
    
    return measurements

if __name__ == "__main__":
    template_path = "/Users/paulius/Downloads/Automatic Slides/output/generated_slides/template.pptx"
    measurements = extract_template_measurements(template_path)
    
    if measurements:
        print("\n=== TEMPLATE MEASUREMENTS ===")
        print(f"Slide size: {measurements['slide_width']} x {measurements['slide_height']}")
        print(f"Found {len(measurements['shapes'])} shapes")
        
        # Save measurements to a file for easy reference
        with open('template_measurements.txt', 'w') as f:
            f.write("TEMPLATE MEASUREMENTS\n")
            f.write("===================\n\n")
            f.write(f"Slide size: {measurements['slide_width']} x {measurements['slide_height']}\n\n")
            
            for shape in measurements['shapes']:
                f.write(f"Shape {shape['index']}: {shape['name']}\n")
                f.write(f"  Type: {shape['type']}\n")
                f.write(f"  Position: Left={shape['left']}, Top={shape['top']}\n")
                f.write(f"  Size: Width={shape['width']}, Height={shape['height']}\n")
                if 'text' in shape:
                    f.write(f"  Text: {shape['text']}\n")
                if 'font_size' in shape:
                    f.write(f"  Font size: {shape['font_size']}\n")
                f.write("\n")
        
        print("\nMeasurements saved to template_measurements.txt")
    else:
        print("Failed to extract measurements") 