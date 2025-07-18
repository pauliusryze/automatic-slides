import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from template_config import TEMPLATE_CONFIG

def set_font(paragraph, size, color, bold=False):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = 'Montserrat'

def add_link_textbox(slide, left, top, width, height, text, url):
    """Create a text box and apply a hyperlink to the text run."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.clear()  # Clear any existing text
    run = p.add_run()
    run.text = text
    run.hyperlink.address = url
    set_font(p, 10, RGBColor(32, 102, 148), bold=True)
    p.alignment = PP_ALIGN.LEFT
    return box

def generate_pptx_from_json_and_images(json_path, image_folder, output_path):
    """Generate PPTX from JSON data and images."""
    
    # Load JSON data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(TEMPLATE_CONFIG['slide_width'])
    prs.slide_height = Inches(TEMPLATE_CONFIG['slide_height'])
    
    for slide_data in data['slides']:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        for shape_config in TEMPLATE_CONFIG['shapes']:
            if shape_config['name'] == 'Title':
                title_box = slide.shapes.add_textbox(
                    Inches(shape_config['left']), 
                    Inches(shape_config['top']), 
                    Inches(shape_config['width']), 
                    Inches(shape_config['height'])
                )
                title_frame = title_box.text_frame
                title_frame.clear()
                # Enable word wrap for title
                title_frame.word_wrap = True
                title_frame.auto_size = None
                
                p = title_frame.paragraphs[0]
                p.text = str(slide_data['title'])
                set_font(p, shape_config['font_size'], RGBColor(153, 81, 44), bold=True)
                p.alignment = PP_ALIGN.LEFT
                
            elif shape_config['name'] == 'Main Image':
                main_image_path = os.path.join(image_folder, slide_data['images']['main'])
                if os.path.exists(main_image_path):
                    slide.shapes.add_picture(
                        main_image_path, 
                        Inches(shape_config['left']), 
                        Inches(shape_config['top']), 
                        Inches(shape_config['width']), 
                        Inches(shape_config['height'])
                    )
                    
            elif shape_config['name'] == 'AA1 Label':
                aa1_box = slide.shapes.add_textbox(
                    Inches(shape_config['left']), 
                    Inches(shape_config['top']), 
                    Inches(shape_config['width']), 
                    Inches(shape_config['height'])
                )
                aa1_frame = aa1_box.text_frame
                aa1_frame.clear()
                p = aa1_frame.paragraphs[0]
                p.text = str(slide_data.get('ad_account', ''))
                set_font(p, shape_config['font_size'], RGBColor(0, 0, 0), bold=True)
                p.alignment = PP_ALIGN.LEFT
                
            elif shape_config['name'] == 'Metrics':
                metrics_box = slide.shapes.add_textbox(
                    Inches(shape_config['left']), 
                    Inches(shape_config['top']), 
                    Inches(shape_config['width']), 
                    Inches(shape_config['height'])
                )
                metrics_frame = metrics_box.text_frame
                metrics_frame.clear()
                metrics_frame.word_wrap = True
                metrics_frame.auto_size = None
                
                # Build metrics text from individual fields
                metrics_text = []
                metrics = slide_data.get('metrics', {})
                
                if metrics.get('spend'):
                    metrics_text.append(f"Spend - ${metrics['spend']:.2f}")
                if metrics.get('cac_scientific'):
                    metrics_text.append(f"CAC (Hyros Scientific) - ${metrics['cac_scientific']:.2f}")
                if metrics.get('cac_last_click'):
                    metrics_text.append(f"CAC (Hyros Last Click) - ${metrics['cac_last_click']:.2f}")
                if metrics.get('cac_clicks_only'):
                    metrics_text.append(f"CAC (Clicks Only - Northbeam) - ${metrics['cac_clicks_only']:.2f}")
                if metrics.get('cac_last_non_direct'):
                    metrics_text.append(f"CAC (Last non-direct touch - Northbeam) - ${metrics['cac_last_non_direct']:.2f}")
                if metrics.get('new_visits_pct'):
                    metrics_text.append(f"New Visits (Northbeam) - {metrics['new_visits_pct']:.2f}%")
                if metrics.get('ecr_pct'):
                    metrics_text.append(f"ECR (Northbeam) - {metrics['ecr_pct']:.2f}%")
                if metrics.get('date_created'):
                    metrics_text.append(f"Date Created - {metrics['date_created']}")
                
                for i, line in enumerate(metrics_text):
                    p = metrics_frame.add_paragraph() if i > 0 else metrics_frame.paragraphs[0]
                    p.text = line
                    color = RGBColor(0, 0, 0)
                    if 'Spend' in line or 'CAC' in line:
                        color = RGBColor(153, 81, 44)
                    elif 'New Visits' in line or 'ECR' in line:
                        color = RGBColor(34, 139, 34)
                    set_font(p, shape_config['font_size'], color)
                    p.alignment = PP_ALIGN.LEFT
                    p.level = 0
                    p.bullet = True
                    
            elif shape_config['name'] == 'Notes':
                # Instead of a single box, create two separate boxes for FB and IG links
                links = slide_data.get('links', {})
                fb_left = shape_config['left']
                fb_top = shape_config['top']
                fb_width = shape_config['width'] / 2.2
                fb_height = shape_config['height']
                ig_left = shape_config['left'] + fb_width + 0.1
                ig_top = shape_config['top']
                ig_width = shape_config['width'] / 2.2
                ig_height = shape_config['height']
                if 'facebook' in links:
                    add_link_textbox(slide, fb_left, fb_top, fb_width, fb_height, links['facebook']['text'], links['facebook']['url'])
                if 'instagram' in links:
                    add_link_textbox(slide, ig_left, ig_top, ig_width, ig_height, links['instagram']['text'], links['instagram']['url'])
                
            elif shape_config['name'] == 'Bottom Image 1':
                bottom_left_image_path = os.path.join(image_folder, slide_data['images']['bottom_left'])
                if os.path.exists(bottom_left_image_path):
                    slide.shapes.add_picture(
                        bottom_left_image_path, 
                        Inches(shape_config['left']), 
                        Inches(shape_config['top']), 
                        Inches(shape_config['width']), 
                        Inches(shape_config['height'])
                    )
                    
            elif shape_config['name'] == 'Bottom Image 2':
                bottom_right_image_path = os.path.join(image_folder, slide_data['images']['bottom_right'])
                if os.path.exists(bottom_right_image_path):
                    slide.shapes.add_picture(
                        bottom_right_image_path, 
                        Inches(shape_config['left']), 
                        Inches(shape_config['top']), 
                        Inches(shape_config['width']), 
                        Inches(shape_config['height'])
                    )
    
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    # Test the JSON-based generator
    json_path = "slide_data.json"
    image_folder = "output/generated_slides"
    output_path = "output/generated_slides/json_test_slides.pptx"
    
    if os.path.exists(json_path):
        result = generate_pptx_from_json_and_images(json_path, image_folder, output_path)
        print(f"✅ Generated slides from JSON: {result}")
    else:
        print(f"❌ JSON file not found: {json_path}") 