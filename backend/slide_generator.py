import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys
import datetime
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from template_config import TEMPLATE_CONFIG

LOG_FILE = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'slidegen_debug.log')
STATIC_UPLOADS = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'static', 'uploads')

# Updated static images mapping with new descriptive names
STATIC_IMAGES = {
    'logo': 'logo.png',
    'mushroom_accent': 'mushroom_accent.png',
    'left_brush_stroke': 'left_brush_stroke.png',
    'right_brush_stroke': 'right_brush_stroke.png',
}


def log(msg):
    with open(LOG_FILE, 'a') as f:
        f.write(f"{datetime.datetime.now().isoformat()} | {msg}\n")


def set_font(paragraph, size, color, bold=False):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = 'Montserrat'


def get_image_path(filename, image_folder):
    if pd.isna(filename) or not filename:
        return None
    # First look in the temp upload dir
    path = os.path.join(image_folder, filename)
    if os.path.exists(path):
        return path
    # Fallback to static/uploads
    static_path = os.path.join(STATIC_UPLOADS, filename)
    if os.path.exists(static_path):
        return static_path
    return None


def generate_pptx_from_csv_and_images(csv_path, image_paths, output_path, image_folder):
    log(f"\n=== GENERATION STARTED ===\nIMAGE_FOLDER: {image_folder}\nCSV: {csv_path}\nOutput: {output_path}")
    df = pd.read_csv(csv_path, encoding='latin1')
    prs = Presentation()
    prs.slide_width = Inches(TEMPLATE_CONFIG['slide_width'])
    prs.slide_height = Inches(TEMPLATE_CONFIG['slide_height'])
    
    for index, row in df.iterrows():
        log(f"[Slide {index+1}] Title: {row.get('title', 'Untitled')}")
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Per-row image logic (use main_image, chart_image, table_image columns)
        per_row_images = {
            'main_creative': get_image_path(row.get('main_image', None), image_folder),
            # SWAP chart and table for placement
            'bottom_chart': get_image_path(row.get('table_image', None), image_folder),
            'bottom_table': get_image_path(row.get('chart_image', None), image_folder),
        }
        for slot, img in per_row_images.items():
            if img:
                log(f"  Found {slot}: {img}")
            else:
                log(f"  MISSING {slot}: {row.get(slot, None)}")

        for shape_config in TEMPLATE_CONFIG['shapes']:
            slot_name = shape_config['name']
            # Handle text shapes (no 'type' or not 'picture')
            if shape_config.get('type', None) != 'picture':
                # Text shape logic
                if slot_name == 'Title':
                    text = str(row.get('title', shape_config.get('text', 'Untitled')))
                elif slot_name == 'AA1 Label':
                    text = str(row.get('ad_account', shape_config.get('text', 'AA1')))
                elif slot_name == 'Notes':
                    fb_text = str(row.get('fb_link_text', 'FB Link'))
                    fb_url = str(row.get('fb_link_url', ''))
                    ig_text = str(row.get('ig_link_text', 'IG Link'))
                    ig_url = str(row.get('ig_link_url', ''))
                    text = f"{fb_text}\n{ig_text}"
                elif slot_name == 'Metrics':
                    text = f"Spend - ${row.get('spend', 0):.2f}\nCAC (Hyros Scientific) - ${row.get('cac_scientific', 0):.2f}\nCAC (Hyros Last Click) - ${row.get('cac_last_click', 0):.2f}\nCAC (Clicks Only - Northbeam) - ${row.get('cac_clicks_only', 0):.2f}\nCAC (Last non-direct touch - Northbeam) - ${row.get('cac_last_non_direct', 0):.2f}\nNew Visits (Northbeam) - {row.get('new_visits_pct', 0):.2f}%\nECR (Northbeam) - {row.get('ecr_pct', 0):.2f}%\nDate Created - {row.get('date_created', 'N/A')}"
                else:
                    text = shape_config.get('text', '')
                log(f"    Placing text shape: {slot_name} with text: {text[:40]}...")
                box = slide.shapes.add_textbox(
                    Inches(shape_config['left']),
                    Inches(shape_config['top']),
                    Inches(shape_config['width']),
                    Inches(shape_config['height'])
                )
                frame = box.text_frame
                frame.clear()
                p = frame.paragraphs[0]
                p.text = text
                set_font(p, shape_config.get('font_size', 12.0), RGBColor(0, 0, 0), bold=(slot_name == 'AA1 Label'))
                p.alignment = PP_ALIGN.LEFT
            elif slot_name in per_row_images:
                img_path = per_row_images[slot_name]
                if img_path:
                    log(f"    Placing {slot_name}: {img_path}")
                    try:
                        slide.shapes.add_picture(
                            img_path,
                            Inches(shape_config['left']),
                            Inches(shape_config['top']),
                            Inches(shape_config['width']),
                            Inches(shape_config['height'])
                        )
                        log(f"    SUCCESS: {slot_name} placed.")
                    except Exception as e:
                        log(f"    ERROR placing {slot_name}: {e}")
                else:
                    log(f"    WARNING: {slot_name} missing for this slide.")
            elif slot_name in STATIC_IMAGES:
                img_path = get_image_path(STATIC_IMAGES[slot_name], image_folder)
                if img_path:
                    log(f"    Placing static {slot_name}: {img_path}")
                    try:
                        slide.shapes.add_picture(
                            img_path,
                            Inches(shape_config['left']),
                            Inches(shape_config['top']),
                            Inches(shape_config['width']),
                            Inches(shape_config['height'])
                        )
                        log(f"    SUCCESS: static {slot_name} placed.")
                    except Exception as e:
                        log(f"    ERROR placing static {slot_name}: {e}")
                else:
                    log(f"    WARNING: static {slot_name} image not found: {STATIC_IMAGES[slot_name]}")
            # else: skip (not an image slot)
    
    prs.save(output_path)
    log(f"=== GENERATION COMPLETE ===\n")
    return output_path
